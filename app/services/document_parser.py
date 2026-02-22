# app/services/document_parser.py
import fitz
from docx import Document
from pptx import Presentation
from io import BytesIO

def extract_text(file_bytes: bytes, content_type: str, filename: str) -> str:
    if content_type == "application/pdf":
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        return "\n".join(page.get_text() for page in doc)

    if "wordprocessingml" in content_type:
        doc = Document(BytesIO(file_bytes))
        return "\n".join(p.text for p in doc.paragraphs)

    if "presentationml" in content_type:
        prs = Presentation(BytesIO(file_bytes))
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    if content_type == "text/plain":
        return file_bytes.decode("utf-8", errors="ignore")

    # 🖼️ JPEG / JPG (OCR)
    if content_type in ("image/jpeg", "image/jpg"):
        image = Image.open(BytesIO(file_bytes))
        text = pytesseract.image_to_string(image)
        return text.strip()

    raise ValueError("Unsupported document type")
